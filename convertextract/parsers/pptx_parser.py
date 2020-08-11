import pptx

from convertextract.parsers.utils import BaseParser

class Parser(BaseParser):
    """Extract text from pptx file using python-pptx
    """

    def extract(self, filename, **kwargs):
        converted_filename = filename[:-5] + '_converted.pptx'
        if 'mapping' in kwargs and kwargs['mapping']:
            transducer = self.create_transducer(kwargs['mapping'])
        else:
            transducer = self.get_transducer(kwargs.get('input_language', ''), kwargs.get('output_language', ''))
        presentation = pptx.Presentation(filename)
        text_runs = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = transducer(run.text).output_string
                        text_runs.append(run.text)
        if "no_write" in kwargs and kwargs['no_write']:
            pass
        else:
            presentation.save(converted_filename)
        return '\n\n'.join(text_runs)